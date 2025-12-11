import io

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    Presentation = None

from typing import List

from sayou.core.registry import register_component

from ..interfaces.base_parser import BaseDocumentParser
from ..models import (
    BaseElement,
    BoundingBox,
    ChartElement,
    Document,
    ElementMetadata,
    ImageElement,
    Slide,
    TableCell,
    TableElement,
    TextElement,
)


@register_component("parser")
class PptxParser(BaseDocumentParser):
    """
    (Tier 2) Parser for Microsoft PowerPoint (.pptx) presentations.

    Iterates through slides, extracting shapes (Text, Picture, Table, Chart).
    Recursively handles grouped shapes and extracts speaker notes.
    """

    component_name = "PptxParser"
    SUPPORTED_TYPES = [".pptx"]

    @classmethod
    def can_handle(cls, file_bytes: bytes, file_name: str) -> float:
        """
        Checks for Zip signature (PK..) or OLE signature.
        """
        # PPTX (Zip)
        if file_bytes.startswith(b"PK\x03\x04") and file_name.lower().endswith(".pptx"):
            return 1.0
        # Legacy DOC
        if file_bytes.startswith(b"\xd0\xcf\x11\xe0"):
            return 1.0
        # Extension fallback
        if any(file_name.lower().endswith(t) for t in cls.SUPPORTED_TYPES):
            return 0.8
        return 0.0

    def _do_parse(self, file_bytes: bytes, file_name: str, **kwargs) -> Document:
        """
        Parse PPTX bytes into a structured Document.

        Args:
            file_bytes (bytes): Binary content of the .pptx file.
            file_name (str): Original filename.
            **kwargs: Options like 'ocr_images'.

        Returns:
            Document: A document object with 'doc_type="slide"'.
        """
        if Presentation is None:
            raise ImportError("python-pptx is required. (pip install python-pptx)")

        try:
            prs = Presentation(io.BytesIO(file_bytes))
        except Exception as e:
            raise ValueError(f"Failed to load PPTX: {e}")

        pages_list = []

        for slide_idx, slide in enumerate(prs.slides):
            elements_list: List[BaseElement] = []
            page_num = slide_idx + 1

            # 1. Shape 순회 (kwargs 전달)
            for shape in slide.shapes:
                elements_list.extend(self._process_shape(shape, page_num, **kwargs))

            # 2. 좌표 기준 정렬
            elements_list.sort(
                key=lambda x: (x.bbox.y0 if x.bbox else 0, x.bbox.x0 if x.bbox else 0)
            )

            # 3. 노트(발표자 메모) 추출
            note_text = ""
            if slide.has_notes_slide:
                note_text = slide.notes_slide.notes_text_frame.text.strip()

            page_text_dump = "\n".join(
                [
                    (
                        e.text
                        if e.type == "text"
                        else (e.text_representation if e.type == "chart" else "")
                    )
                    for e in elements_list
                ]
            )

            if note_text:
                page_text_dump += "\n[NOTE]: " + note_text

            page_obj = Slide(
                page_num=page_num,
                elements=elements_list,
                text=page_text_dump,
                note_text=note_text if note_text else None,
                has_notes=bool(note_text),
            )
            pages_list.append(page_obj)

        return Document(
            file_name=file_name,
            file_id=file_name,
            doc_type="slide",
            page_count=len(pages_list),
            pages=pages_list,
        )

    def _process_shape(self, shape, page_num: int, **kwargs) -> List[BaseElement]:
        """
        Convert a PPTX shape into one or more Sayou Elements.

        Handles recursive extraction for Group shapes and type-specific
        logic for TextFrames, Pictures, Tables, and Charts.

        Args:
            shape (Shape): The python-pptx Shape object.
            page_num (int): Slide number.
            **kwargs: Configuration passed from parse.

        Returns:
            List[BaseElement]: A list of extracted elements (Text, Image, Table, Chart).
        """
        extracted = []

        # 1. 그룹 (Group) -> 재귀 호출
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                extracted.extend(self._process_shape(sub_shape, page_num, **kwargs))
            return extracted

        # 공통 BBox
        try:
            bbox = BoundingBox(
                x0=shape.left,
                y0=shape.top,
                x1=shape.left + shape.width,
                y1=shape.top + shape.height,
            )
        except AttributeError:
            bbox = None

        meta = ElementMetadata(
            page_num=page_num, id=f"p{page_num}:shape:{shape.shape_id}"
        )
        ocr_enabled = kwargs.get("ocr_images", True)
        placeholder_type = None
        if shape.is_placeholder:
            placeholder_type = str(shape.placeholder_format.type)

        # 2. 텍스트 (Text Frame)
        if shape.has_text_frame and shape.text.strip():
            text_elem = TextElement(
                id=meta.id,
                type="text",
                bbox=bbox,
                meta=meta,
                text=shape.text.strip(),
                raw_attributes={
                    "shape_type": str(shape.shape_type),
                    "placeholder_type": placeholder_type,
                },
            )
            extracted.append(text_elem)

        # 3. 이미지 (Picture)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                img_elem = self._process_image_data(
                    image_bytes=shape.image.blob,
                    img_format=shape.image.ext,
                    elem_id=meta.id,
                    page_num=meta.page_num,
                    bbox=bbox,
                    ocr_enabled=ocr_enabled,
                )
                img_elem.raw_attributes["placeholder_type"] = placeholder_type
                extracted.append(img_elem)
            except Exception as e:
                self._log(f"PPT Image Error: {e}", level="warning")

        # 4. 테이블 (Table)
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            try:
                table = shape.table
                rows_data = []
                rows_cells = []  # TODO: v0.1.0+ - TableCell 상세 구현

                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    rows_data.append(row_data)

                table_elem = TableElement(
                    id=meta.id,
                    type="table",
                    bbox=bbox,
                    meta=meta,
                    data=rows_data,
                    cells=[],  # v0.1.0+ 구현
                    raw_attributes={"placeholder_type": placeholder_type},
                )
                extracted.append(table_elem)
            except Exception as e:
                self._log(f"PPT Table Error: {e}", level="warning")

        # 5. 차트 (Chart)
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            try:
                chart = shape.chart
                chart_title = (
                    chart.chart_title.text_frame.text if chart.has_title else "Chart"
                )
                chart_type_str = str(chart.chart_type)

                # 차트 데이터를 텍스트로 단순 변환
                text_rep = f"Chart: {chart_title} (Type: {chart_type_str})\n"
                for i, series in enumerate(chart.series):
                    series_name = series.name or f"Series {i+1}"
                    series_data = ", ".join([str(v) for v in (series.values or [])])
                    text_rep += f"- {series_name}: [{series_data}]\n"

                chart_elem = ChartElement(
                    id=meta.id,
                    type="chart",
                    bbox=bbox,
                    meta=meta,
                    chart_title=chart_title,
                    chart_type=chart_type_str,
                    text_representation=text_rep,
                    raw_attributes={
                        "placeholder_type": placeholder_type,
                        "category_count": len(chart.categories),
                        "series_count": len(chart.series),
                    },
                )
                extracted.append(chart_elem)
            except Exception as e:
                self._log(f"PPT Chart Error: {e}", level="warning")

        return extracted
