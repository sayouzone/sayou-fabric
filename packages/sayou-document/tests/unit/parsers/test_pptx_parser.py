"""
Unit tests for PptxParser.

Verifies:
- can_handle scoring
- Slide.text field populated (Bug #2 regression)
- Speaker notes extracted
- Shape types: text, picture, table, chart
"""

from __future__ import annotations

from unittest.mock import MagicMock, PropertyMock, patch

import pytest

from sayou.document.models import (ChartElement, Document, ImageElement, Slide,
                                   TableElement, TextElement)
from sayou.document.parser.pptx_parser import PptxParser

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _make_text_shape(text: str, slide_num: int = 1, shape_id: int = 1):
    shape = MagicMock()
    shape.shape_id = shape_id
    shape.has_text_frame = True
    shape.text = text
    shape.left = 100
    shape.top = 100
    shape.width = 400
    shape.height = 50
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    shape.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
    shape.is_placeholder = False
    return shape


def _make_picture_shape(shape_id: int = 2):
    shape = MagicMock()
    shape.shape_id = shape_id
    shape.has_text_frame = False
    shape.text = ""
    shape.left = 200
    shape.top = 200
    shape.width = 100
    shape.height = 100
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    shape.shape_type = MSO_SHAPE_TYPE.PICTURE
    shape.is_placeholder = False
    shape.image.blob = b"\x89PNG\r\n\x1a\n" + b"\x00" * 8
    shape.image.ext = "png"
    return shape


def _make_prs(slides_data: list[dict]):
    """Build a mock Presentation with configurable slides."""
    prs = MagicMock()
    mock_slides = []
    for slide_data in slides_data:
        slide = MagicMock()
        slide.shapes = slide_data.get("shapes", [])
        slide.has_notes_slide = bool(slide_data.get("note"))
        if slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = slide_data["note"]
        mock_slides.append(slide)
    prs.slides = mock_slides
    return prs


# ---------------------------------------------------------------------------
# can_handle
# ---------------------------------------------------------------------------


class TestPptxParserCanHandle:
    def test_pk_with_pptx_extension(self):
        assert PptxParser.can_handle(b"PK\x03\x04", "deck.pptx") == 1.0

    def test_ole_signature(self):
        assert PptxParser.can_handle(b"\xd0\xcf\x11\xe0", "old.ppt") == 1.0

    def test_extension_fallback(self):
        assert PptxParser.can_handle(b"\x00\x00", "deck.pptx") == 0.8

    def test_non_pptx_returns_0(self):
        assert PptxParser.can_handle(b"%PDF", "report.pdf") == 0.0


# ---------------------------------------------------------------------------
# _do_parse
# ---------------------------------------------------------------------------


class TestPptxParserParse:
    @patch("sayou.document.parser.pptx_parser.Presentation")
    def test_returns_document(self, MockPrs):
        MockPrs.return_value = _make_prs([{"shapes": []}])
        parser = PptxParser()
        doc = parser._do_parse(b"PK\x03\x04", "deck.pptx")
        assert isinstance(doc, Document)
        assert doc.doc_type == "slide"

    @patch("sayou.document.parser.pptx_parser.Presentation")
    def test_slide_count(self, MockPrs):
        MockPrs.return_value = _make_prs([{"shapes": []}, {"shapes": []}])
        parser = PptxParser()
        doc = parser._do_parse(b"PK", "deck.pptx")
        assert doc.page_count == 2
        assert len(doc.pages) == 2

    @patch("sayou.document.parser.pptx_parser.Presentation")
    def test_slide_text_field_populated(self, MockPrs):
        """Regression: Slide.text must not be silently dropped."""
        shapes = [_make_text_shape("Slide title text")]
        MockPrs.return_value = _make_prs([{"shapes": shapes}])
        parser = PptxParser()
        doc = parser._do_parse(b"PK", "deck.pptx")
        slide = doc.pages[0]
        assert isinstance(slide, Slide)
        assert slide.text is not None
        assert "Slide title text" in slide.text

    @patch("sayou.document.parser.pptx_parser.Presentation")
    def test_speaker_notes_extracted(self, MockPrs):
        shapes = [_make_text_shape("Main content")]
        MockPrs.return_value = _make_prs(
            [{"shapes": shapes, "note": "Remember to pause here."}]
        )
        parser = PptxParser()
        doc = parser._do_parse(b"PK", "deck.pptx")
        slide = doc.pages[0]
        assert slide.has_notes is True
        assert slide.note_text == "Remember to pause here."
        assert "[NOTE]: Remember to pause here." in slide.text

    @patch("sayou.document.parser.pptx_parser.Presentation")
    def test_text_shape_produces_text_element(self, MockPrs):
        shapes = [_make_text_shape("Hello from shape")]
        MockPrs.return_value = _make_prs([{"shapes": shapes}])
        parser = PptxParser()
        doc = parser._do_parse(b"PK", "deck.pptx")
        text_elems = [e for e in doc.pages[0].elements if isinstance(e, TextElement)]
        assert len(text_elems) >= 1
        assert any(e.text == "Hello from shape" for e in text_elems)

    @patch("sayou.document.parser.pptx_parser.Presentation")
    def test_picture_shape_produces_image_element(self, MockPrs):
        shapes = [_make_picture_shape()]
        MockPrs.return_value = _make_prs([{"shapes": shapes}])
        parser = PptxParser()
        doc = parser._do_parse(b"PK", "deck.pptx")
        img_elems = [e for e in doc.pages[0].elements if isinstance(e, ImageElement)]
        assert len(img_elems) == 1
        assert img_elems[0].image_base64 is not None

    @patch("sayou.document.parser.pptx_parser.Presentation")
    def test_table_shape_produces_table_element(self, MockPrs):
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        shape = MagicMock()
        shape.shape_id = 3
        shape.has_text_frame = False
        shape.text = ""
        shape.left = shape.top = shape.width = shape.height = 0
        shape.shape_type = MSO_SHAPE_TYPE.TABLE
        shape.is_placeholder = False

        mock_row = MagicMock()
        mock_row.cells = [MagicMock(text="A"), MagicMock(text="B")]
        shape.table.rows = [mock_row]

        MockPrs.return_value = _make_prs([{"shapes": [shape]}])
        parser = PptxParser()
        doc = parser._do_parse(b"PK", "deck.pptx")
        tbl_elems = [e for e in doc.pages[0].elements if isinstance(e, TableElement)]
        assert len(tbl_elems) == 1
        assert tbl_elems[0].data[0] == ["A", "B"]
